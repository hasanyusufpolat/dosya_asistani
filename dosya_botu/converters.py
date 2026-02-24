"""
PROFESYONEL DOSYA DÖNÜŞTÜRME MODÜLÜ
Tüm dönüşümler yüksek kalitede ve sorunsuz çalışır
Gelişmiş tipografi, tablo yönetimi ve format koruma
"""

import os
import logging
import datetime
from PIL import Image, ImageEnhance, ImageFilter
import pytesseract

# Tesseract yolunu ayarla (Windows için)
pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

# Loglama ayarları
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('converters.log'),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)

# ========== YARDIMCI FONKSİYONLAR ==========
def clean_text(text):
    """Metni temizle ve düzenle"""
    if not text:
        return ""
    # Fazla boşlukları temizle
    lines = text.split('\n')
    cleaned_lines = [' '.join(line.split()) for line in lines if line.strip()]
    return '\n'.join(cleaned_lines)

def detect_table_structure(text):
    """Metin içinde tablo yapısını tespit et"""
    lines = text.split('\n')
    tablo_olasiligi = 0
    
    for line in lines:
        # Birden fazla boşluk veya sekme varsa tablo olabilir
        if '\t' in line or '  ' in line:
            tablo_olasiligi += 1
        # Düzenli sütunlar varsa
        if '|' in line:
            tablo_olasiligi += 2
    
    return tablo_olasiligi > len(lines) * 0.3

def format_number(value):
    """Sayıları formatla"""
    if isinstance(value, (int, float)):
        if value.is_integer():
            return str(int(value))
        else:
            return f"{value:.2f}".replace('.', ',')
    return str(value)

# ========== WORD DÖNÜŞÜMLERİ (PROFESYONEL) ==========
def word_to_pdf(input_path, output_path):
    """Word -> PDF (PROFESYONEL - TİPOGRAFİ KORUMALI)"""
    try:
        from docx import Document
        from docx.shared import RGBColor as DocxRGB
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.units import cm
        from reportlab.lib import colors
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont
        
        doc = Document(input_path)
        c = canvas.Canvas(output_path, pagesize=A4)
        width, height = A4
        
        # Sayfa kenar boşlukları
        left_margin = 2*cm
        right_margin = width - 2*cm
        y = height - 2*cm
        line_height = 0.6*cm
        
        # Başlık stilleri için font boyutları
        title_size = 16
        heading1_size = 14
        heading2_size = 13
        normal_size = 11
        
        for paragraph in doc.paragraphs:
            if not paragraph.text.strip():
                y -= line_height * 0.5
                continue
            
            # Paragraf stilini belirle
            style_name = paragraph.style.name.lower() if paragraph.style else "normal"
            font_size = normal_size
            is_bold = False
            is_italic = False
            
            if 'title' in style_name:
                font_size = title_size
                is_bold = True
            elif 'heading 1' in style_name:
                font_size = heading1_size
                is_bold = True
            elif 'heading 2' in style_name:
                font_size = heading2_size
                is_bold = True
            
            # Run bazlı stilleri kontrol et
            for run in paragraph.runs:
                if run.bold:
                    is_bold = True
                if run.italic:
                    is_italic = True
            
            # Font stilini ayarla
            if is_bold and is_italic:
                c.setFont("Helvetica-BoldOblique", font_size)
            elif is_bold:
                c.setFont("Helvetica-Bold", font_size)
            elif is_italic:
                c.setFont("Helvetica-Oblique", font_size)
            else:
                c.setFont("Helvetica", font_size)
            
            # Metni kelimelere böl ve satırlara ayır
            words = paragraph.text.split()
            current_line = ""
            
            for word in words:
                test_line = current_line + " " + word if current_line else word
                # Karakter genişliği yaklaşık font_size * 0.6
                if len(test_line) * (font_size * 0.6) < (right_margin - left_margin):
                    current_line = test_line
                else:
                    if y < line_height + 1*cm:
                        c.showPage()
                        y = height - 2*cm
                        # Fontu yeniden ayarla
                        if is_bold and is_italic:
                            c.setFont("Helvetica-BoldOblique", font_size)
                        elif is_bold:
                            c.setFont("Helvetica-Bold", font_size)
                        elif is_italic:
                            c.setFont("Helvetica-Oblique", font_size)
                        else:
                            c.setFont("Helvetica", font_size)
                    
                    c.drawString(left_margin, y, current_line)
                    y -= line_height
                    current_line = word
            
            # Kalan satırı yaz
            if current_line:
                if y < line_height + 1*cm:
                    c.showPage()
                    y = height - 2*cm
                    # Fontu yeniden ayarla
                    if is_bold and is_italic:
                        c.setFont("Helvetica-BoldOblique", font_size)
                    elif is_bold:
                        c.setFont("Helvetica-Bold", font_size)
                    elif is_italic:
                        c.setFont("Helvetica-Oblique", font_size)
                    else:
                        c.setFont("Helvetica", font_size)
                
                c.drawString(left_margin, y, current_line)
                y -= line_height * 1.2
        
        c.save()
        logger.info(f"✅ Word -> PDF dönüşüm başarılı: {input_path}")
        return True, None
    except Exception as e:
        logger.error(f"❌ Word -> PDF dönüşüm hatası: {e}")
        return False, f"Word -> PDF dönüşüm hatası: {str(e)}"

def word_to_excel(input_path, output_path):
    """Word -> Excel (PROFESYONEL - AKILLI TABLO ALGILAMA)"""
    try:
        from docx import Document
        import pandas as pd
        from openpyxl import Workbook
        from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
        from openpyxl.utils import get_column_letter
        
        doc = Document(input_path)
        
        # Önce tabloları dene
        tables_data = []
        for table in doc.tables:
            table_rows = []
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                table_rows.append(row_data)
            if table_rows:
                tables_data.append(table_rows)
        
        # Tablo varsa onları kullan
        if tables_data:
            # En büyük tabloyu seç
            main_table = max(tables_data, key=len)
            df = pd.DataFrame(main_table[1:], columns=main_table[0] if main_table[0] else None)
        else:
            # Tablo yoksa paragrafları topla
            data = []
            for para in doc.paragraphs:
                if para.text.strip():
                    # Tablo yapısı var mı kontrol et
                    if detect_table_structure(para.text):
                        # Tablo satırlarını ayır
                        lines = para.text.split('\n')
                        for line in lines:
                            if '\t' in line:
                                cells = line.split('\t')
                                data.append([c.strip() for c in cells])
                            elif '  ' in line:
                                cells = [c for c in line.split('  ') if c.strip()]
                                data.append(cells)
                            else:
                                data.append([line.strip()])
                    else:
                        data.append([para.text.strip()])
            
            df = pd.DataFrame(data, columns=['İçerik'] if data else None)
        
        # PROFESYONEL EXCEL OLUŞTUR
        with pd.ExcelWriter(output_path, engine='openpyxl') as writer:
            df.to_excel(writer, index=False, sheet_name='Dönüştürülen Veri')
            
            # Excel stil ayarları
            workbook = writer.book
            worksheet = writer.sheets['Dönüştürülen Veri']
            
            # Başlık stili
            header_font = Font(name='Calibri', size=12, bold=True, color="FFFFFF")
            header_fill = PatternFill(start_color="2E75B6", end_color="2E75B6", fill_type="solid")
            header_alignment = Alignment(horizontal="center", vertical="center")
            
            # Hücre kenarlıkları
            thin_border = Border(
                left=Side(style='thin'),
                right=Side(style='thin'),
                top=Side(style='thin'),
                bottom=Side(style='thin')
            )
            
            # Başlık satırını formatla
            for cell in worksheet[1]:
                cell.font = header_font
                cell.fill = header_fill
                cell.alignment = header_alignment
                cell.border = thin_border
            
            # Veri hücrelerini formatla
            for row in worksheet.iter_rows(min_row=2, max_row=len(df)+1):
                for cell in row:
                    cell.font = Font(name='Calibri', size=11)
                    cell.alignment = Alignment(horizontal="left", vertical="center", wrap_text=True)
                    cell.border = thin_border
                    
                    # Sayısal değerleri otomatik formatla
                    if isinstance(cell.value, (int, float)):
                        cell.number_format = '#,##0.00'
            
            # Sütun genişliklerini otomatik ayarla
            for column in worksheet.columns:
                max_length = 0
                column_letter = get_column_letter(column[0].column)
                for cell in column:
                    try:
                        if cell.value and len(str(cell.value)) > max_length:
                            max_length = len(str(cell.value))
                    except:
                        pass
                adjusted_width = min(max_length + 2, 50)
                worksheet.column_dimensions[column_letter].width = adjusted_width
            
            # Satır yüksekliklerini ayarla
            worksheet.row_dimensions[1].height = 25
            for i in range(2, len(df)+2):
                worksheet.row_dimensions[i].height = 18
        
        logger.info(f"✅ Word -> Excel dönüşüm başarılı: {input_path}")
        return True, None
    except Exception as e:
        logger.error(f"❌ Word -> Excel dönüşüm hatası: {e}")
        return False, f"Word -> Excel dönüşüm hatası: {str(e)}"

def word_to_pptx(input_path, output_path):
    """Word -> PowerPoint (PROFESYONEL - TASARIM ODAKLI)"""
    try:
        from docx import Document
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        
        doc = Document(input_path)
        prs = Presentation()
        
        # Slayt tasarımı
        title_slide_layout = prs.slide_layouts[0]
        content_slide_layout = prs.slide_layouts[1]
        
        # Ana başlık slaytı
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        if title:
            title.text = "WORD DÖKÜMANI DÖNÜŞÜMÜ"
            title.text_frame.paragraphs[0].font.size = Pt(48)
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
            title.text_frame.paragraphs[0].font.bold = True
        
        if subtitle:
            subtitle.text = f"Kaynak: {os.path.basename(input_path)}\nTarih: {datetime.datetime.now().strftime('%d.%m.%Y')}"
            subtitle.text_frame.paragraphs[0].font.size = Pt(20)
        
        # İçerik slaytları
        content_slides = []
        current_slide = None
        current_text_frame = None
        items_per_slide = 6
        
        # Önce tüm paragrafları topla
        all_paragraphs = []
        for para in doc.paragraphs:
            if para.text.strip():
                # Başlık mı içerik mi kontrol et
                style_name = para.style.name.lower() if para.style else "normal"
                is_heading = 'heading' in style_name or 'title' in style_name
                all_paragraphs.append({
                    'text': para.text.strip(),
                    'is_heading': is_heading,
                    'style': style_name
                })
        
        # Slaytları oluştur
        slide_count = 0
        for i, para in enumerate(all_paragraphs):
            if para['is_heading'] or i % items_per_slide == 0:
                # Yeni slayt oluştur
                slide_count += 1
                slide = prs.slides.add_slide(content_slide_layout)
                
                # Başlık
                title = slide.shapes.title
                if title:
                    if para['is_heading']:
                        title.text = para['text'][:50]
                    else:
                        title.text = f"İçerik - Sayfa {slide_count}"
                    title.text_frame.paragraphs[0].font.size = Pt(32)
                    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
                    title.text_frame.paragraphs[0].font.bold = True
                
                # İçerik alanı
                content = slide.placeholders[1]
                text_frame = content.text_frame
                text_frame.clear()
                
                if not para['is_heading']:
                    p = text_frame.add_paragraph()
                    p.text = para['text']
                    p.font.size = Pt(20)
                    p.font.color.rgb = RGBColor(0, 0, 0)
                    p.level = 0
                    if p.text:
                        p.text = "• " + p.text
                
                current_text_frame = text_frame
            else:
                # Mevcut slayta ekle
                if current_text_frame:
                    p = current_text_frame.add_paragraph()
                    p.text = para['text']
                    p.font.size = Pt(20)
                    p.font.color.rgb = RGBColor(0, 0, 0)
                    p.level = 1
                    if p.text:
                        p.text = "  • " + p.text
        
        prs.save(output_path)
        logger.info(f"✅ Word -> PowerPoint dönüşüm başarılı: {input_path}")
        return True, None
    except Exception as e:
        logger.error(f"❌ Word -> PowerPoint dönüşüm hatası: {e}")
        return False, f"Word -> PowerPoint dönüşüm hatası: {str(e)}"

# ========== EXCEL DÖNÜŞÜMLERİ (PROFESYONEL) ==========
def excel_to_pdf(input_path, output_path):
    """Excel -> PDF (PROFESYONEL - TABLO KORUMALI)"""
    try:
        import pandas as pd
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import A4, landscape
        from reportlab.lib.units import cm
        from reportlab.platypus import Table, TableStyle
        from reportlab.lib import colors
        
        df = pd.read_excel(input_path)
        
        # Veriyi temizle
        df = df.fillna('')
        
        # PDF'i yatay formatda oluştur
        c = canvas.Canvas(output_path, pagesize=landscape(A4))
        width, height = landscape(A4)
        
        # Verileri hazırla
        data = [df.columns.tolist()] + df.values.tolist()
        
        # Hücre genişliklerini hesapla (içeriğe göre)
        col_widths = []
        for i, col in enumerate(df.columns):
            max_len = len(str(col))
            for val in df.iloc[:, i]:
                if pd.notna(val):
                    max_len = max(max_len, len(str(val)))
            # Karakter başına 0.3 cm
            col_widths.append(min(max_len * 0.25*cm + 1*cm, 6*cm))
        
        # PROFESYONEL TABLO OLUŞTUR
        table = Table(data, colWidths=col_widths, repeatRows=1)
        table.setStyle(TableStyle([
            # Başlık satırı
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#2E75B6')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 10),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 8),
            ('TOPPADDING', (0, 0), (-1, 0), 8),
            
            # Veri satırları
            ('BACKGROUND', (0, 1), (-1, -1), colors.HexColor('#F2F2F2')),
            ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor('#666666')),
            ('FONTSIZE', (0, 1), (-1, -1), 9),
            ('ROWBACKGROUNDS', (0, 1), (-1, -1), 
             [colors.HexColor('#FFFFFF'), colors.HexColor('#F2F2F2')]),
            
            # Sayısal değerleri sağa hizala
            ('ALIGN', (1, 1), (-1, -1), 'RIGHT'),
        ]))
        
        # Sayfa boyutunu hesapla
        table_height = len(data) * 0.5*cm + 2*cm
        if table_height > height - 3*cm:
            # Çoklu sayfa desteği
            from reportlab.platypus import SimpleDocTemplate, PageBreak
            from reportlab.lib.pagesizes import A4
            import tempfile
            
            # Geçici PDF oluştur
            with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp:
                doc = SimpleDocTemplate(tmp.name, pagesize=landscape(A4))
                elements = [table]
                doc.build(elements)
                
                # Geçici dosyayı hedefe kopyala
                import shutil
                shutil.copy2(tmp.name, output_path)
                os.unlink(tmp.name)
        else:
            # Tek sayfa
            table.wrapOn(c, width, height)
            table.drawOn(c, 1*cm, height - table_height - 2*cm)
        
        # Alt bilgi ekle
        c.setFont("Helvetica", 8)
        c.setFillColor(colors.HexColor('#666666'))
        c.drawString(1*cm, 1*cm, f"Oluşturulma: {datetime.datetime.now().strftime('%d.%m.%Y %H:%M')}")
        c.drawRightString(width - 1*cm, 1*cm, f"Sayfa 1 / 1")
        
        c.save()
        logger.info(f"✅ Excel -> PDF dönüşüm başarılı: {input_path}")
        return True, None
    except Exception as e:
        logger.error(f"❌ Excel -> PDF dönüşüm hatası: {e}")
        return False, f"Excel -> PDF dönüşüm hatası: {str(e)}"

def excel_to_word(input_path, output_path):
    """Excel -> Word (PROFESYONEL - SAYFAYA TAM SIĞDIRMA - DÜZELTİLDİ)"""
    try:
        import pandas as pd
        from docx import Document
        from docx.shared import Inches, Cm, Pt, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.enum.table import WD_TABLE_ALIGNMENT
        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement
        
        # Excel'i oku - tüm sayfaları al
        excel_file = pd.ExcelFile(input_path)
        sheet_names = excel_file.sheet_names
        
        doc = Document()
        
        # Sayfa yapısı ayarları - A4 boyutu için optimize
        section = doc.sections[0]
        section.page_width = Cm(21)  # A4 genişlik
        section.page_height = Cm(29.7)  # A4 yükseklik
        section.left_margin = Cm(1.5)
        section.right_margin = Cm(1.5)
        section.top_margin = Cm(1.5)
        section.bottom_margin = Cm(1.5)
        
        # Ana başlık - şık tasarım
        title = doc.add_heading('EXCEL DÖKÜMANI DÖNÜŞÜMÜ', 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = title.runs[0]
        run.font.size = Pt(24)
        run.font.bold = True
        run.font.name = 'Calibri'
        run.font.color.rgb = RGBColor(0, 51, 102)
        
        # Alt başlık bilgileri
        doc.add_paragraph()
        info_para = doc.add_paragraph()
        info_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = info_para.add_run(f"📊 Kaynak: {os.path.basename(input_path)}")
        run.font.size = Pt(11)
        run.font.name = 'Calibri'
        run.font.color.rgb = RGBColor(100, 100, 100)
        
        info_para2 = doc.add_paragraph()
        info_para2.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = info_para2.add_run(f"📅 Dönüşüm: {datetime.datetime.now().strftime('%d.%m.%Y %H:%M')}")
        run.font.size = Pt(11)
        run.font.name = 'Calibri'
        run.font.color.rgb = RGBColor(100, 100, 100)
        
        # Çizgi ekle
        doc.add_paragraph('_' * 80)
        doc.add_paragraph()
        
        # Toplam sayfa sayısı
        doc.add_paragraph(f"📑 Toplam {len(sheet_names)} sayfa", style='Intense Quote')
        doc.add_paragraph()
        
        # Her Excel sayfası için
        for sheet_idx, sheet_name in enumerate(sheet_names):
            df = pd.read_excel(input_path, sheet_name=sheet_name)
            df = df.fillna('')  # NaN'leri boş string yap
            
            # Sayfa başlığı
            if sheet_idx > 0:
                doc.add_page_break()
            
            heading = doc.add_heading(f'Sayfa {sheet_idx + 1}: {sheet_name}', level=1)
            for run in heading.runs:
                run.font.size = Pt(18)
                run.font.color.rgb = RGBColor(0, 102, 204)
            
            if df.empty:
                doc.add_paragraph("📭 Bu sayfa boş")
                doc.add_paragraph()
                continue
            
            # Kullanılabilir genişlik
            available_width_cm = 18  # Yaklaşık kullanılabilir genişlik
            
            # Sütun genişliklerini hesapla - sayfaya tam sığdır
            col_count = len(df.columns)
            
            # Minimum genişlik (karakter başına 0.15 cm, minimum 2 cm)
            min_widths = []
            for i, col in enumerate(df.columns):
                max_len = len(str(col))
                for val in df.iloc[:, i]:
                    if val != '':
                        max_len = max(max_len, len(str(val)))
                # Her karakter yaklaşık 0.15 cm, minimum 2 cm
                min_width = max(max_len * 0.15, 2)
                min_widths.append(min_width)
            
            # Toplam minimum genişlik
            total_min_width = sum(min_widths)
            
            # Eğer toplam minimum genişlik sayfadan büyükse, orantılı olarak küçült
            if total_min_width > available_width_cm:
                # Orantılı olarak küçült
                scale_factor = available_width_cm / total_min_width
                col_widths = [w * scale_factor for w in min_widths]
            else:
                # Fazla boşluk varsa, eşit dağıt
                extra_space = available_width_cm - total_min_width
                extra_per_col = extra_space / col_count
                col_widths = [w + extra_per_col for w in min_widths]
            
            # TABLO OLUŞTUR - TEK PARÇA (sayfalama yok)
            rows, cols = df.shape
            table = doc.add_table(rows=rows+1, cols=cols)
            table.style = 'Light Grid Accent 1'
            table.alignment = WD_TABLE_ALIGNMENT.CENTER
            table.autofit = False
            
            # Sütun genişliklerini ayarla
            for i, width in enumerate(col_widths):
                for row in table.rows:
                    row.cells[i].width = Cm(width)
            
            # BAŞLIK SATIRI
            for col in range(cols):
                cell = table.cell(0, col)
                cell.text = str(df.columns[col])
                
                # Hücre arka plan rengi (koyu mavi)
                tc = cell._tc
                tcPr = tc.get_or_add_tcPr()
                shd = OxmlElement('w:shd')
                shd.set(qn('w:fill'), '2E75B6')
                tcPr.append(shd)
                
                # Metin formatı
                for paragraph in cell.paragraphs:
                    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    run = paragraph.runs[0]
                    run.font.bold = True
                    run.font.size = Pt(11)
                    run.font.name = 'Calibri'
                    run.font.color.rgb = RGBColor(255, 255, 255)
            
            # VERİ SATIRLARI
            for row in range(rows):
                for col in range(cols):
                    cell = table.cell(row+1, col)
                    value = df.iloc[row, col]
                    
                    # Değeri formatla
                    if isinstance(value, (int, float)):
                        if isinstance(value, float) and not value.is_integer():
                            cell.text = f"{value:.2f}".replace('.', ',')
                        else:
                            cell.text = str(int(value) if isinstance(value, float) else value)
                    else:
                        cell.text = str(value)
                    
                    # Hücre arka planı (alternatif renkler)
                    if row % 2 == 0:
                        tc = cell._tc
                        tcPr = tc.get_or_add_tcPr()
                        shd = OxmlElement('w:shd')
                        shd.set(qn('w:fill'), 'F2F2F2')
                        tcPr.append(shd)
                    
                    # Metin formatı
                    for paragraph in cell.paragraphs:
                        # Sayısal değerleri sağa, metinleri sola hizala
                        if isinstance(df.iloc[row, col], (int, float)):
                            paragraph.alignment = WD_ALIGN_PARAGRAPH.RIGHT
                        else:
                            paragraph.alignment = WD_ALIGN_PARAGRAPH.LEFT
                        
                        run = paragraph.runs[0]
                        run.font.size = Pt(10)
                        run.font.name = 'Calibri'
        
        doc.save(output_path)
        logger.info(f"✅ Excel -> Word dönüşüm başarılı: {input_path}")
        return True, None
    except Exception as e:
        logger.error(f"❌ Excel -> Word dönüşüm hatası: {e}")
        import traceback
        traceback.print_exc()
        return False, f"Excel -> Word dönüşüm hatası: {str(e)}"

def excel_to_pptx(input_path, output_path):
    """Excel -> PowerPoint (PROFESYONEL - GRAFİK DESTEKLİ)"""
    try:
        import pandas as pd
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE
        from pptx.chart.data import ChartData
        
        df = pd.read_excel(input_path)
        df = df.fillna('')
        
        prs = Presentation()
        
        rows, cols = df.shape
        rows_per_slide = 18  # Her slayta 18 satır
        
        # Ana başlık slaytı
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        if title:
            title.text = "EXCEL VERİLERİ"
            title.text_frame.paragraphs[0].font.size = Pt(48)
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
            title.text_frame.paragraphs[0].font.bold = True
        
        if subtitle:
            subtitle.text = f"Toplam {rows} satır, {cols} sütun\n{os.path.basename(input_path)}"
            subtitle.text_frame.paragraphs[0].font.size = Pt(20)
        
        # Veri slaytları
        for slide_start in range(0, rows, rows_per_slide):
            slide_end = min(slide_start + rows_per_slide, rows)
            
            # Yeni slayt
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Başlık ekle
            title_box = slide.shapes.add_textbox(
                int(Inches(0.5).emu), 
                int(Inches(0.2).emu), 
                int(Inches(9).emu), 
                int(Inches(0.8).emu)
            )
            title_frame = title_box.text_frame
            title_frame.text = f"Excel Verileri - Sayfa {slide_start//rows_per_slide + 1}"
            title_frame.paragraphs[0].font.size = Pt(32)
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
            
            # Verileri hazırla
            table_data = [df.columns.tolist()] + df.iloc[slide_start:slide_end].values.tolist()
            
            rows_in_slide = len(table_data)
            cols_in_slide = len(table_data[0])
            
            # Tablo boyutları
            left = int(Inches(0.5).emu)
            top = int(Inches(1.5).emu)
            width = int(Inches(9).emu)
            height = int(Inches(5.0).emu)
            
            # Tabloyu oluştur
            table = slide.shapes.add_table(rows_in_slide, cols_in_slide, left, top, width, height).table
            
            # Sütun genişliklerini ayarla
            col_width = int(width / cols_in_slide)
            for col in range(cols_in_slide):
                table.columns[col].width = col_width
            
            # Verileri doldur ve formatla
            for row in range(rows_in_slide):
                for col in range(cols_in_slide):
                    cell = table.cell(row, col)
                    cell.text = str(table_data[row][col])
                    
                    # Hücre formatı
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.size = Pt(12)
                        paragraph.alignment = PP_ALIGN.CENTER
                    
                    # Başlık satırı formatı
                    if row == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(46, 117, 182)
                        for paragraph in cell.text_frame.paragraphs:
                            paragraph.font.color.rgb = RGBColor(255, 255, 255)
                            paragraph.font.bold = True
                            paragraph.font.size = Pt(13)
                    else:
                        # Alternatif satır renkleri
                        if row % 2 == 0:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = RGBColor(242, 242, 242)
        
        # Grafik slaytı (eğer sayısal veri varsa)
        numeric_cols = df.select_dtypes(include=['number']).columns
        if len(numeric_cols) >= 2 and len(df) <= 20:
            chart_slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(chart_slide_layout)
            
            title = slide.shapes.title
            if title:
                title.text = "Veri Grafiği"
            
            # Grafik verilerini hazırla
            chart_data = ChartData()
            chart_data.categories = df.iloc[:10, 0].astype(str).tolist()
            
            for col in numeric_cols[1:3]:  # İlk 2 sayısal sütun
                chart_data.add_series(str(col), df.iloc[:10][col].tolist())
            
            # Grafik ekle
            left = Inches(1)
            top = Inches(1.5)
            width = Inches(8)
            height = Inches(5)
            
            slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data
            )
        
        prs.save(output_path)
        logger.info(f"✅ Excel -> PowerPoint dönüşüm başarılı: {input_path}")
        return True, None
    except Exception as e:
        logger.error(f"❌ Excel -> PowerPoint dönüşüm hatası: {e}")
        return False, f"Excel -> PowerPoint dönüşüm hatası: {str(e)}"

# ========== POWERPOINT DÖNÜŞÜMLERİ (PROFESYONEL) ==========
def pptx_to_pdf(input_path, output_path):
    """PowerPoint -> PDF (PROFESYONEL - TASARIM KORUMALI)"""
    try:
        from pptx import Presentation
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import A4, landscape
        from reportlab.lib.units import cm
        from reportlab.lib import colors
        
        prs = Presentation(input_path)
        c = canvas.Canvas(output_path, pagesize=landscape(A4))
        width, height = landscape(A4)
        
        left_margin = 2*cm
        y = height - 2*cm
        line_height = 0.7*cm
        
        for slide_num, slide in enumerate(prs.slides, 1):
            # Slayt başlığı
            c.setFont("Helvetica-Bold", 18)
            c.setFillColor(colors.HexColor('#2E75B6'))
            c.drawString(left_margin, y, f"Slayt {slide_num}")
            y -= line_height * 2
            
            c.setFont("Helvetica", 11)
            c.setFillColor(colors.black)
            
            # Slayt içeriği
            shapes_with_text = 0
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    shapes_with_text += 1
                    
                    # Metni düzenle
                    text = shape.text.strip()
                    
                    # Başlık mı içerik mi kontrol et
                    if len(text) < 50 and not '\n' in text:
                        c.setFont("Helvetica-Bold", 14)
                        c.setFillColor(colors.HexColor('#1E4E7C'))
                    else:
                        c.setFont("Helvetica", 11)
                        c.setFillColor(colors.black)
                    
                    # Metni satırlara böl
                    text_lines = text.split('\n')
                    
                    for line in text_lines:
                        if line.strip():
                            if y < line_height + 1*cm:
                                c.showPage()
                                y = height - 2*cm
                                c.setFont("Helvetica", 11)
                                c.setFillColor(colors.black)
                            
                            # Madde işareti ekle
                            c.drawString(left_margin + 0.5*cm, y, f"• {line.strip()}")
                            y -= line_height
                    
                    y -= line_height * 0.5
            
            # Eğer içerik yoksa
            if shapes_with_text == 0:
                c.setFont("Helvetica-Oblique", 11)
                c.setFillColor(colors.HexColor('#666666'))
                c.drawString(left_margin + 0.5*cm, y, "(Bu slaytta metin yok)")
                y -= line_height
            
            # Sonraki slayt için yeni sayfa
            if slide_num < len(prs.slides):
                c.showPage()
                y = height - 2*cm
        
        c.save()
        logger.info(f"✅ PowerPoint -> PDF dönüşüm başarılı: {input_path}")
        return True, None
    except Exception as e:
        logger.error(f"❌ PowerPoint -> PDF dönüşüm hatası: {e}")
        return False, f"PowerPoint -> PDF dönüşüm hatası: {str(e)}"

def pptx_to_word(input_path, output_path):
    """PowerPoint -> Word (PROFESYONEL)"""
    try:
        from pptx import Presentation
        from docx import Document
        from docx.shared import Inches, Pt
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        
        prs = Presentation(input_path)
        doc = Document()
        
        # Stil ayarları
        style = doc.styles['Normal']
        style.font.name = 'Arial'
        style.font.size = Pt(11)
        
        # Ana başlık
        title = doc.add_heading('POWERPOINT DÖKÜMANI DÖNÜŞÜMÜ', 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        title.runs[0].font.size = Pt(24)
        
        doc.add_paragraph(f"Kaynak dosya: {os.path.basename(input_path)}")
        doc.add_paragraph(f"Toplam {len(prs.slides)} slayt")
        doc.add_paragraph(f"Dönüşüm tarihi: {datetime.datetime.now().strftime('%d.%m.%Y %H:%M')}")
        doc.add_paragraph()
        
        for slide_num, slide in enumerate(prs.slides, 1):
            # Slayt başlığı
            doc.add_heading(f'Slayt {slide_num}', level=1)
            
            # Slayt içeriği
            shapes_with_text = 0
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    shapes_with_text += 1
                    
                    # Metni temizle ve düzenle
                    text = shape.text.strip()
                    
                    # Başlık mı içerik mi kontrol et
                    if len(text) < 50 and not '\n' in text:
                        # Alt başlık olarak ekle
                        doc.add_heading(text, level=2)
                    else:
                        # Normal paragraf olarak ekle
                        paragraphs = text.split('\n')
                        for para in paragraphs:
                            if para.strip():
                                p = doc.add_paragraph()
                                p.add_run(para.strip()).font.size = Pt(11)
                                p.paragraph_format.left_indent = Inches(0.3)
                                p.paragraph_format.space_after = Pt(3)
            
            # Eğer içerik yoksa
            if shapes_with_text == 0:
                doc.add_paragraph("(Bu slaytta metin yok)")
            
            # Sayfa sonu (son slayt hariç)
            if slide_num < len(prs.slides):
                doc.add_page_break()
        
        doc.save(output_path)
        logger.info(f"✅ PowerPoint -> Word dönüşüm başarılı: {input_path}")
        return True, None
    except Exception as e:
        logger.error(f"❌ PowerPoint -> Word dönüşüm hatası: {e}")
        return False, f"PowerPoint -> Word dönüşüm hatası: {str(e)}"

# ========== PDF DÖNÜŞÜMLERİ (PROFESYONEL) ==========
def pdf_to_word(input_path, output_path):
    """PDF -> Word (PROFESYONEL - METİN KORUMALI)"""
    try:
        from docx import Document
        from docx.shared import Inches, Pt
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        import PyPDF2
        
        doc = Document()
        
        # Stil ayarları
        style = doc.styles['Normal']
        style.font.name = 'Times New Roman'
        style.font.size = Pt(12)
        
        with open(input_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            
            # Ana başlık
            title = doc.add_heading('PDF DÖKÜMANI DÖNÜŞÜMÜ', 0)
            title.alignment = WD_ALIGN_PARAGRAPH.CENTER
            title.runs[0].font.size = Pt(24)
            title.runs[0].font.name = 'Arial'
            
            doc.add_paragraph(f"Kaynak dosya: {os.path.basename(input_path)}")
            doc.add_paragraph(f"Toplam {len(pdf_reader.pages)} sayfa")
            doc.add_paragraph(f"Dönüşüm tarihi: {datetime.datetime.now().strftime('%d.%m.%Y %H:%M')}")
            doc.add_paragraph()
            
            for page_num in range(len(pdf_reader.pages)):
                page = pdf_reader.pages[page_num]
                text = page.extract_text()
                
                # Sayfa başlığı
                doc.add_heading(f'Sayfa {page_num + 1}', level=1)
                
                if text and text.strip():
                    # Metni temizle ve düzenle
                    lines = text.split('\n')
                    paragraph_text = ""
                    
                    for line in lines:
                        clean_line = ' '.join(line.split())
                        if clean_line:
                            # Paragraf sonu kontrolü (nokta ile bitiyorsa)
                            if clean_line.endswith('.') or clean_line.endswith('!') or clean_line.endswith('?'):
                                paragraph_text += " " + clean_line if paragraph_text else clean_line
                                doc.add_paragraph(paragraph_text)
                                paragraph_text = ""
                            else:
                                if paragraph_text:
                                    paragraph_text += " " + clean_line
                                else:
                                    paragraph_text = clean_line
                    
                    # Kalan metni ekle
                    if paragraph_text:
                        doc.add_paragraph(paragraph_text)
                else:
                    doc.add_paragraph("(Bu sayfada metin bulunamadı)")
                
                # Sayfa sonu (son sayfa hariç)
                if page_num < len(pdf_reader.pages) - 1:
                    doc.add_page_break()
        
        doc.save(output_path)
        logger.info(f"✅ PDF -> Word dönüşüm başarılı: {input_path}")
        return True, None
    except Exception as e:
        logger.error(f"❌ PDF -> Word dönüşüm hatası: {e}")
        return False, f"PDF -> Word dönüşüm hatası: {str(e)}"

# ========== GÖRSEL DÖNÜŞÜMLERİ (SADECE GÖRSEL -> PDF ve GÖRSEL -> WORD) ==========
def image_to_pdf(input_path, output_path):
    """Görsel -> PDF (PROFESYONEL - YÜKSEK KALİTE)"""
    try:
        from PIL import Image
        import img2pdf
        
        # Görseli yükle ve optimize et
        image = Image.open(input_path)
        
        # Görseli RGB'ye çevir
        if image.mode != 'RGB':
            image = image.convert('RGB')
        
        # Görseli yeniden boyutlandır (A4 sığacak şekilde)
        a4_width, a4_height = 2480, 3508  # 300 DPI
        image.thumbnail((a4_width, a4_height), Image.Resampling.LANCZOS)
        
        # Geçici olarak yüksek kalitede kaydet
        temp_image_path = input_path + "_temp.jpg"
        image.save(temp_image_path, 'JPEG', quality=100, optimize=True, dpi=(300,300))
        
        # PDF'e çevir
        with open(temp_image_path, "rb") as f:
            image_data = f.read()
        
        pdf_bytes = img2pdf.convert(temp_image_path)
        
        with open(output_path, "wb") as f:
            f.write(pdf_bytes)
        
        # Geçici dosyayı temizle
        os.remove(temp_image_path)
        
        logger.info(f"✅ Görsel -> PDF dönüşüm başarılı: {input_path}")
        return True, None
    except Exception as e:
        logger.error(f"❌ Görsel -> PDF dönüşüm hatası: {e}")
        return False, f"Görsel -> PDF dönüşüm hatası: {str(e)}"

def image_to_word(input_path, output_path):
    """Görsel -> Word (OCR - PROFESYONEL - YÜKSEK KALİTE - DÜZELTİLDİ)"""
    try:
        from PIL import Image, ImageEnhance, ImageFilter
        from docx import Document
        from docx.shared import Inches, Pt, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        import pytesseract
        
        # Görseli yükle ve ön işle
        image = Image.open(input_path)
        
        # Görseli büyüt (OCR kalitesi için)
        width, height = image.size
        if width < 2000:
            new_size = (width * 2, height * 2)
            image = image.resize(new_size, Image.Resampling.LANCZOS)
        
        # Görseli gri tonlamaya çevir
        image = image.convert('L')
        
        # Kontrast artır
        enhancer = ImageEnhance.Contrast(image)
        image = enhancer.enhance(2.5)
        
        # Gürültü azalt
        image = image.filter(ImageFilter.MedianFilter(size=3))
        
        # Kenar keskinleştir
        image = image.filter(ImageFilter.SHARPEN)
        
        # Geçici olarak kaydet
        temp_image_path = input_path + "_temp_ocr.png"
        image.save(temp_image_path, 'PNG', dpi=(300,300))
        
        # OCR ile metin çıkar (Türkçe ve İngilizce)
        custom_config = r'--oem 3 --psm 6 -l tur+eng'
        text = pytesseract.image_to_string(temp_image_path, config=custom_config)
        
        # Word belgesi oluştur - profesyonel tasarım
        doc = Document()
        
        # Sayfa yapısı ayarları
        section = doc.sections[0]
        section.left_margin = Inches(1)
        section.right_margin = Inches(1)
        section.top_margin = Inches(1)
        section.bottom_margin = Inches(1)
        
        # Stil ayarları
        style = doc.styles['Normal']
        style.font.name = 'Calibri'
        style.font.size = Pt(11)
        
        # Başlık - şık tasarım
        title = doc.add_heading('📄 GÖRSELDEN OCR İLE DÖNÜŞTÜRÜLEN METİN', 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = title.runs[0]
        run.font.size = Pt(26)
        run.font.bold = True
        run.font.name = 'Calibri'
        run.font.color.rgb = RGBColor(0, 51, 102)
        
        # Alt başlık bilgileri
        doc.add_paragraph()
        info_para = doc.add_paragraph()
        info_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = info_para.add_run(f"🖼️ Kaynak: {os.path.basename(input_path)}")
        run.font.size = Pt(12)
        run.font.color.rgb = RGBColor(100, 100, 100)
        
        info_para2 = doc.add_paragraph()
        info_para2.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = info_para2.add_run(f"📅 Dönüşüm: {datetime.datetime.now().strftime('%d.%m.%Y %H:%M')}")
        run.font.size = Pt(12)
        run.font.color.rgb = RGBColor(100, 100, 100)
        
        # Çizgi ekle
        doc.add_paragraph('_' * 80)
        doc.add_paragraph()
        
        # Metin içeriği
        if text.strip():
            # Metni paragraflara böl ve temizle
            raw_paragraphs = text.split('\n\n')
            
            for para in raw_paragraphs:
                if para.strip():
                    # Satırları birleştir
                    lines = para.split('\n')
                    clean_text = ' '.join([line.strip() for line in lines if line.strip()])
                    
                    # Paragraf ekle
                    p = doc.add_paragraph()
                    run = p.add_run(clean_text)
                    run.font.size = Pt(11)
                    run.font.name = 'Calibri'
                    
                    # Paragraf arası boşluk
                    p.paragraph_format.space_after = Pt(12)
        else:
            p = doc.add_paragraph()
            run = p.add_run("(Görselde metin bulunamadı veya okunamadı)")
            run.font.size = Pt(12)
            run.font.italic = True
            run.font.color.rgb = RGBColor(150, 150, 150)
        
        # Sayfa sonu ekle
        doc.add_page_break()
        
        # Orijinal görseli ekle - yeni sayfada
        doc.add_heading('🖼️ ORİJİNAL GÖRSEL', level=1)
        doc.add_paragraph()
        
        # Görseli ekle (orta hizalı)
        try:
            doc.add_picture(input_path, width=Inches(5))
            last_paragraph = doc.paragraphs[-1]
            last_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
        except:
            doc.add_paragraph("(Görsel yüklenemedi)")
        
        # Geçici dosyayı temizle
        os.remove(temp_image_path)
        
        doc.save(output_path)
        logger.info(f"✅ Görsel -> Word dönüşüm başarılı: {input_path}")
        return True, None
    except Exception as e:
        logger.error(f"❌ Görsel -> Word dönüşüm hatası: {e}")
        import traceback
        traceback.print_exc()
        return False, f"Görsel -> Word dönüşüm hatası: {str(e)}"

# ========== DÖNÜŞTÜRME YÖNETİCİSİ ==========
async def convert_file(input_path, output_path, source_type, target_type):
    """
    Tüm dönüşümleri yöneten ana fonksiyon
    """
    
    conversion_functions = {
        # Word dönüşümleri
        ('WORD', 'PDF'): word_to_pdf,
        ('WORD', 'EXCEL'): word_to_excel,
        ('WORD', 'POWERPOINT'): word_to_pptx,
        
        # Excel dönüşümleri
        ('EXCEL', 'PDF'): excel_to_pdf,
        ('EXCEL', 'WORD'): excel_to_word,
        ('EXCEL', 'POWERPOINT'): excel_to_pptx,
        
        # PowerPoint dönüşümleri
        ('POWERPOINT', 'PDF'): pptx_to_pdf,
        ('POWERPOINT', 'WORD'): pptx_to_word,
        
        # PDF dönüşümleri
        ('PDF', 'WORD'): pdf_to_word,
        
        # Görsel dönüşümleri (SADECE bunlar kaldı)
        ('GORSEL', 'PDF'): image_to_pdf,
        ('GORSEL', 'WORD'): image_to_word,
    }
    
    func = conversion_functions.get((source_type, target_type))
    if func:
        try:
            logger.info(f"🔄 Dönüşüm başlıyor: {source_type} -> {target_type}")
            import asyncio
            result, error = await asyncio.get_event_loop().run_in_executor(
                None, func, input_path, output_path
            )
            if result:
                logger.info(f"✅ Dönüşüm tamamlandı: {source_type} -> {target_type}")
            else:
                logger.error(f"❌ Dönüşüm başarısız: {source_type} -> {target_type} - {error}")
            return result, error
        except Exception as e:
            logger.error(f"❌ Dönüşüm hatası: {e}")
            return False, str(e)
    else:
        logger.warning(f"⚠️ Desteklenmeyen dönüşüm: {source_type} -> {target_type}")
        return False, f"Desteklenmeyen dönüşüm: {source_type} -> {target_type}"